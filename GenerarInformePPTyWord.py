# Autor: Roy Ayala Galvis
# uso: crear informes mensuales del estado de los servidores
# Pre requsitos: tienen q estar abierto el navegador con las paginas que se quiere capturar obligatoriamente

import pyautogui
import time
from PIL import Image, ImageChops
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PPTXInches

# Coordenadas de la región a capturar (puedes ajustarlas según tus necesidades)
REGION = (40, 100, 1400, 800)  # (left, top, width, height)

# Array con los títulos de las imágenes - 26 titulos por que se definio estaticamente 26 captura de tab de chrome
IMAGE_TITLES = [
    "Servidores Blade HP DESA1",
    "Servidores Blade HP DESA2",
    "Servidores Blade HP DESA3",
    "Servidores Blade IBM Testing1",
    "Servidores Blade IBM Testing2",
    "Servidores Blade IBM Testing3",
    "Servidores Blade CISCO Stage1",
    "Servidores Blade CISCO Stage2",
    "Servidores Blade CISCO Stage3",
    "Servidores Blade IBM de Prod1",
    "Servidores Blade IBM de Prod2",
    "Servidores Blade IBM de Prod3",
    "Subscripcion de Azure Costos",
    "Subscripcion de AWS Costos",
    "Estado de Balanceador 1",
    "Estado de Balanceador 2",
    "Dashboard de Grafana 1",
    "Dashboard de Grafana 2",
    "Dashboard Dynatrace 1",
    "Dashboard Dynatrace 2",
    "Dashboard Wazuh 1 windows",
    "Dashboard Wazuh 2 linux",
    "Dashboard HaProxy Testing",
    "Dashboard HaProxy PROD",
    "Dashboard Darktrace ",
    "Dashboar de Red"
    # Añade más títulos según sea necesario
]

def capture_chrome_tabs():
    # Activa la ventana de Chrome
    pyautogui.hotkey('alt', 'tab')
    time.sleep(2)  # Aumentar el tiempo para asegurar que Chrome esté activo

    screenshots = []
    tab_count = 0
    first_tab = None

    while True:
        # Captura la ventana activa
        window = pyautogui.getActiveWindow()
        if window is None:
            print("No se pudo encontrar la ventana activa.")
            break

        # Captura la ventana
        screenshot = pyautogui.screenshot(region=REGION)
        filename = f"tab_{tab_count}.png"
        screenshot.save(filename)
        screenshots.append(filename)

        # Si es la primera pestaña, guardamos la imagen
        if tab_count == 0:
            first_tab = screenshot

        # Intenta cambiar a la siguiente pestaña
        pyautogui.hotkey('ctrl', 'tab')
        time.sleep(2)  # Aumentar el tiempo para permitir que la pestaña cambie

        # Comprueba si hemos vuelto a la primera pestaña
        if tab_count > 0:
            current_tab = pyautogui.screenshot(region=REGION)
            diff = ImageChops.difference(first_tab, current_tab)
            if diff.getbbox() is None:
                print("Se ha vuelto a la primera pestaña. Terminando la captura.")
                break

        tab_count += 1

        # Limitar el número de pestañas capturadas para evitar bucles infinitos
        if tab_count >= 26:  # Cambia este número según la cantidad de pestañas que esperas capturar
            print("Se ha alcanzado el límite de captura de pestañas.")
            break

    return screenshots

def create_word_document(screenshots):
    doc = Document()

    # Página de Carátula
    doc.add_heading('Informe Mensual de estado de Servidores', level=1)
    doc.add_paragraph('Estado actual de Servidores de produccion del ultimo mes')
    doc.add_paragraph('Autor: Raimundo Ayala Galvis')
    doc.add_paragraph('Especialista en Servidores - SRE')
    doc.add_paragraph('Fecha: 23 de agosto de 2024')
    doc.add_page_break()  # Salto de página

    # Página de Índice de Contenido
    doc.add_heading('Mes de Proceso', level=1)
    doc.add_paragraph('Agosto 2024')

    doc.add_page_break()  # Salto de página

    # Capturas de Pantalla
    doc.add_heading('Capturas de pantalla de Chrome', level=1)

    for i, screenshot in enumerate(screenshots):
        title = IMAGE_TITLES[i] if i < len(IMAGE_TITLES) else f'Imagen {i+1}'
        doc.add_heading(title, level=2)
        doc.add_picture(screenshot, width=Inches(6))
        doc.add_paragraph(f'Imagen: {i+1}', style='Caption')

    # Página de Conclusiones y Recomendaciones
    doc.add_page_break()  # Salto de página
    doc.add_heading('Conclusiones y Recomendaciones', level=1)
    doc.add_paragraph('Aquí puedes añadir tus conclusiones y recomendaciones sobre el contenido capturado.')

    # Añadir tabla de contenido
    doc.add_page_break()  # Salto de página
    doc.add_heading('Tabla de Contenido', level=1)
    doc.add_paragraph('Tabla de Contenido dinámica:')
    doc.add_paragraph('1. Capturas de Pantalla de Chrome', style='ListBullet')
    for i in range(len(screenshots)):
        doc.add_paragraph(f'   {i + 1}. Imagen {i + 1}', style='ListBullet2')
    doc.add_paragraph('2. Conclusiones y Recomendaciones', style='ListBullet')

    # Guarda el documento
    doc.save('InformeMensualEstadoServidoresProduccion.docx')

def create_powerpoint_presentation(screenshots):
    prs = Presentation()

    # Diapositiva de Carátula
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Layout de título
    title = slide.shapes.title
    subtitle = slide.placeholders[1]  # Normalmente el segundo placeholder es el subtítulo

    # Establecer el título y el subtítulo
    if title:
        title.text = "Informe Mensual de estado de Servidores"
    if subtitle:
        subtitle.text = "Estado actual de servidores de produccion"

    # Diapositivas de Capturas de Pantalla
    for i, screenshot in enumerate(screenshots):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Diseño de diapositiva en blanco
        
        # Agregar un encabezado manualmente
        left = Inches(1)
        top = Inches(0.5)
        width = Inches(8)
        height = Inches(0.5)
        header = slide.shapes.add_textbox(left, top, width, height)
        text_frame = header.text_frame
        p = text_frame.add_paragraph()

        title = IMAGE_TITLES[i] if i < len(IMAGE_TITLES) else f'Imagen {i+1}'
        
        p.text = f"Informe Mensual - {title}"
        p.font.bold = True  # Puedes personalizar el formato del texto aquí

        # Agregar la imagen de la captura de pantalla
        # Ajustar la posición de la imagen para que no se superponga con el encabezado
        img_left = Inches(1)
        img_top = Inches(1.5)  # Aumentar el valor para mover la imagen más abajo
        img_width = Inches(8)
        img_height = Inches(5)
        img = slide.shapes.add_picture(screenshot, img_left, img_top, width=img_width, height=img_height)

    # Guarda la presentación de PowerPoint
    prs.save('InformeMensualEstadoServidoresProduccion.pptx')

# Captura las pantallas
screenshots = capture_chrome_tabs()

# Crea el documento de Word
if screenshots:  # Verifica que se hayan capturado imágenes
    create_word_document(screenshots)
    create_powerpoint_presentation(screenshots)
    print("Proceso completado. El documento 'InformeMensualEstadoServidoresProduccion.docx' y la presentación 'InformeMensualEstadoServidoresProduccion.pptx' han sido creados.")
else:
    print("No se capturaron imágenes.")